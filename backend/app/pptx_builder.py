"""
Simple PPTX builder utilities that map a DeckModel (JSON/dict) into a python-pptx Presentation

This module provides small helper functions for each slide type. It aims to produce editable text boxes
and simple, clean layouts suitable for job-application decks.

Limitations:
- Font substitution depends on the viewer machine. We attempt to set requested font names, but PowerPoint
  will fall back if not installed.
- Images are downloaded synchronously using requests. Large images may slow generation.
"""

import collections
# Compatibility shim for python-pptx on some Python versions where collections.Container/Mapping
# were moved to collections.abc. This ensures older pptx packages that import collections.Container
# still find the attributes.
try:
    if not hasattr(collections, 'Container'):
        import collections.abc as _c_abc
        collections.Container = _c_abc.Container
    if not hasattr(collections, 'Mapping'):
        import collections.abc as _c_abc2
        collections.Mapping = _c_abc2.Mapping
except Exception:
    pass

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.dml import MSO_THEME_COLOR
from bs4 import BeautifulSoup
import requests
import io
import re
from typing import Dict, Any, List


def _hex_to_rgb(hexcolor: str):
    if not hexcolor:
        return (255, 255, 255)
    h = hexcolor.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


def _is_dark(hexcolor: str) -> bool:
    try:
        r, g, b = _hex_to_rgb(hexcolor or '#FFFFFF')
        # relative luminance
        luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
        return luminance < 0.5
    except Exception:
        return False


def _safe_font(slide, name: str):
    # python-pptx will use given name; there's no easy way to check availability here.
    return name or 'Arial'


def _download_image_to_stream(url: str):
    try:
        resp = requests.get(url, timeout=10)
        resp.raise_for_status()
        return io.BytesIO(resp.content)
    except Exception:
        return None


def set_slide_background(prs, slide, template: Dict[str, Any]):
    # python-pptx lacks full gradient background support; add a full-size rectangle as background.
    width = prs.slide_width
    height = prs.slide_height
    bgType = template.get('bgType', 'solid')
    if bgType == 'image':
        # use bgImage field from template if present, otherwise skip
        img_url = template.get('bgImage') or template.get('bgImageUrl')
        if img_url:
            stream = _download_image_to_stream(img_url)
            if stream:
                try:
                    slide.shapes.add_picture(stream, 0, 0, width, height)
                    return
                except Exception:
                    pass
        # fallback to solid if image missing
        bgType = 'solid'
    elif bgType == 'gradient':
        # approximate: add two rectangles with translucency — simple fallback
        from pptx.util import Pt
        g = template.get('bgGradient', {}) or {}
        from_color = _hex_to_rgb(g.get('from', template.get('bgColor') or '#FFFFFF'))
        to_color = _hex_to_rgb(g.get('to', template.get('bgColor') or '#FFFFFF'))
        # Add full-size shape filled with from_color
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, height)
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*from_color)
        shape.shape_style = 0
        # overlay with semi-transparent rectangle for to_color
        shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, height)
        fill2 = shape2.fill
        fill2.solid()
        fill2.fore_color.rgb = RGBColor(*to_color)
        fill2.transparency = 0.85
        shape2.shape_style = 0
        # send both to back by doing nothing (they are added beneath content)
    else:
        bg = template.get('bgColor')
        if bg:
            rgb = _hex_to_rgb(bg)
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*rgb)
            shape.shape_style = 0


def _apply_title_style_to_paragraph(p, template: Dict[str, Any], is_on_dark_bg: bool = False):
    title_font = template.get('fontTitle') or template.get('fontBody') or 'Arial'
    p.font.name = title_font
    p.font.size = Pt(template.get('titleFontSize', 40))
    p.font.bold = True
    accent = template.get('accentColor')
    if accent and not is_on_dark_bg:
        r, g, b = _hex_to_rgb(accent)
        p.font.color.rgb = RGBColor(r, g, b)
    elif is_on_dark_bg:
        p.font.color.rgb = RGBColor(255, 255, 255)


def create_title_slide(prs: Presentation, slide_def: Dict[str, Any], template: Dict[str, Any]):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    set_slide_background(prs, slide, template)

    title_font = _safe_font(slide, template.get('fontTitle'))
    body_font = _safe_font(slide, template.get('fontBody'))

    title_text = slide_def.get('title') or ''
    subtitle = slide_def.get('subtitle') or ''
    images = slide_def.get('images') or []

    hint = template.get('layoutHints', {}).get('titleSlide', 'centeredBig')

    if hint == 'leftPhotoRightText' and images:
        # left photo area
        left = Inches(0.5)
        top = Inches(1)
        pic_w = Inches(2.2)
        pic_h = Inches(2.8)
        img_stream = _download_image_to_stream(images[0])
        if img_stream:
            slide.shapes.add_picture(img_stream, left, top, width=pic_w, height=pic_h)
        # right text
        tx_left = Inches(3.0)
        tx_top = Inches(1.0)
        tx_w = prs.slide_width - tx_left - Inches(0.5)
        tx_h = Inches(3.5)
        txBox = slide.shapes.add_textbox(tx_left, tx_top, tx_w, tx_h)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        _apply_title_style_to_paragraph(p, template, is_on_dark_bg=_is_dark(template.get('bgColor')))
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.name = body_font
            p2.font.size = Pt(template.get('subtitleFontSize', 18))
    elif hint == 'centeredBig':
        tx_left = Inches(1)
        tx_top = Inches(1.5)
        tx_w = prs.slide_width - Inches(2)
        tx_h = Inches(4)
        txBox = slide.shapes.add_textbox(tx_left, tx_top, tx_w, tx_h)
        tf = txBox.text_frame
        tf.vertical_anchor = None
        p = tf.paragraphs[0]
        p.text = title_text
        _apply_title_style_to_paragraph(p, template, is_on_dark_bg=_is_dark(template.get('bgColor')))
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.name = body_font
            p2.font.size = Pt(template.get('subtitleFontSize', 18))
            p2.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    elif hint == 'fullBleedImage' and images:
        # full-bleed image
        img_stream = _download_image_to_stream(images[0])
        if img_stream:
            slide.shapes.add_picture(img_stream, 0, 0, prs.slide_width, prs.slide_height)
        # overlay rectangle for text
        left = Inches(0.5)
        top = prs.slide_height - Inches(2.2)
        width = prs.slide_width - Inches(1)
        height = Inches(1.8)
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        rect.fill.solid()
        accent = _hex_to_rgb(template.get('accentColor') or '#000000')
        rect.fill.fore_color.rgb = RGBColor(*accent)
        rect.fill.fore_color.brightness = -0.2
        rect.fill.transparency = 0.3
        tf = rect.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title_text
        _apply_title_style_to_paragraph(p, template, is_on_dark_bg=True)
        p.font.color.rgb = RGBColor(255, 255, 255)
    else:
        # default centered
        tx_left = Inches(1)
        tx_top = Inches(1.5)
        tx_w = prs.slide_width - Inches(2)
        tx_h = Inches(3)
        txBox = slide.shapes.add_textbox(tx_left, tx_top, tx_w, tx_h)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        _apply_title_style_to_paragraph(p, template, is_on_dark_bg=_is_dark(template.get('bgColor')))
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.name = body_font
            p2.font.size = Pt(template.get('subtitleFontSize', 18))

    # add contact line at bottom if present in subtitle or notes
    notes = slide_def.get('notes')
    if notes:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = notes

    return slide


def create_summary_slide(prs: Presentation, slide_def: Dict[str, Any], template: Dict[str, Any]):
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    set_slide_background(prs, slide, template)
    title_font = _safe_font(slide, template.get('fontTitle'))
    body_font = _safe_font(slide, template.get('fontBody'))

    # Add a heading strip with accent color
    accent = template.get('accentColor')
    if accent:
        accent_rgb = _hex_to_rgb(accent)
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.3), prs.slide_width, Inches(0.6))
        strip.fill.solid()
        strip.fill.fore_color.rgb = RGBColor(*accent_rgb)
        strip.fill.transparency = 0
        strip.line.fill.background()
        # title on strip
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), prs.slide_width - Inches(1), Inches(0.6)).text_frame
        p = tx.paragraphs[0]
        p.text = slide_def.get('title') or ''
        _apply_title_style_to_paragraph(p, template, is_on_dark_bg=_is_dark(accent))
    # Body area
    tx_left = Inches(0.7)
    tx_top = Inches(1.2)
    tx_w = prs.slide_width - Inches(1.4)
    tx_h = prs.slide_height - tx_top - Inches(1)
    txBox = slide.shapes.add_textbox(tx_left, tx_top, tx_w, tx_h)
    tf = txBox.text_frame
    tf.word_wrap = True
    bullets = slide_def.get('bullets') or []
    for b in bullets:
        p = tf.add_paragraph()
        p.text = "• " + b
        p.font.name = body_font
        p.font.size = Pt(template.get('bodyFontSize', 16))
    notes = slide_def.get('notes')
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def create_experience_slide(prs: Presentation, slide_def: Dict[str, Any], template: Dict[str, Any]):
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    set_slide_background(prs, slide, template)
    title_font = _safe_font(slide, template.get('fontTitle'))
    body_font = _safe_font(slide, template.get('fontBody'))

    title = slide_def.get('title') or 'Experience'
    tx_left = Inches(0.5)
    tx_top = Inches(0.5)
    tx_w = prs.slide_width - Inches(1)
    tx_h = Inches(1)
    txBox = slide.shapes.add_textbox(tx_left, tx_top, tx_w, tx_h)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    _apply_title_style_to_paragraph(p, template, is_on_dark_bg=_is_dark(template.get('bgColor')))

    items = slide_def.get('items') or []
    # render each experience item as heading + bullets
    y = Inches(1.4)
    left = Inches(0.6)
    for item in items:
        comp = item.get('company', '')
        role = item.get('role', '')
        dates = item.get('dates', '')
        head = f"{role} — {comp} ({dates})" if role or comp or dates else ''
        box = slide.shapes.add_textbox(left, y, prs.slide_width - Inches(1.2), Inches(0.6))
        tf2 = box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = head
        p2.font.name = title_font
        p2.font.size = Pt(template.get('bodyFontSize', 16))
        p2.font.bold = True
        y = y + Inches(0.35)
        bullets = item.get('bullets') or []
        for b in bullets:
            bbox = slide.shapes.add_textbox(left + Inches(0.15), y, prs.slide_width - Inches(1.5), Inches(0.4))
            btf = bbox.text_frame
            bp = btf.paragraphs[0]
            bp.text = '• ' + b
            bp.font.name = body_font
            bp.font.size = Pt(template.get('bodyFontSize', 14))
            y = y + Inches(0.35)
        y = y + Inches(0.1)
        if y > prs.slide_height - Inches(1):
            # create new slide if overflow
            slide = prs.slides.add_slide(layout)
            y = Inches(0.5)
    notes = slide_def.get('notes')
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def create_skills_slide(prs: Presentation, slide_def: Dict[str, Any], template: Dict[str, Any]):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    set_slide_background(prs, slide, template)
    title_font = _safe_font(slide, template.get('fontTitle'))
    body_font = _safe_font(slide, template.get('fontBody'))

    title = slide_def.get('title') or 'Skills'
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), prs.slide_width - Inches(1), Inches(0.6)).text_frame
    tx.paragraphs[0].text = title
    tx.paragraphs[0].font.name = title_font
    tx.paragraphs[0].font.size = Pt(template.get('headingFontSize', 22))

    bullets = slide_def.get('bullets') or []
    # split into two columns
    mid = (len(bullets) + 1) // 2
    left_col = bullets[:mid]
    right_col = bullets[mid:]

    left_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(4.5), Inches(4))
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.0), Inches(4))
    ltf = left_box.text_frame
    rtf = right_box.text_frame
    for b in left_col:
        p = ltf.add_paragraph()
        p.text = '• ' + b
        p.font.name = body_font
        p.font.size = Pt(template.get('bodyFontSize', 14))
    for b in right_col:
        p = rtf.add_paragraph()
        p.text = '• ' + b
        p.font.name = body_font
        p.font.size = Pt(template.get('bodyFontSize', 14))

    notes = slide_def.get('notes')
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def create_projects_slide(prs: Presentation, slide_def: Dict[str, Any], template: Dict[str, Any]):
    layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    set_slide_background(prs, slide, template)
    title_font = _safe_font(slide, template.get('fontTitle'))
    body_font = _safe_font(slide, template.get('fontBody'))

    title = slide_def.get('title') or 'Projects'
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), prs.slide_width - Inches(1), Inches(0.6)).text_frame
    tx.paragraphs[0].text = title
    tx.paragraphs[0].font.name = title_font
    tx.paragraphs[0].font.size = Pt(template.get('headingFontSize', 22))

    items = slide_def.get('items') or []
    y = Inches(1.0)
    for it in items:
        t = it.get('title') or ''
        desc = it.get('description') or ''
        img = it.get('image') or None
        # left image if exists
        if img:
            s = _download_image_to_stream(img)
            if s:
                slide.shapes.add_picture(s, Inches(0.5), y, width=Inches(2.5))
                text_left = Inches(3.2)
            else:
                text_left = Inches(0.6)
        else:
            text_left = Inches(0.6)
        box = slide.shapes.add_textbox(text_left, y, prs.slide_width - text_left - Inches(0.5), Inches(1.2))
        btf = box.text_frame
        btf.paragraphs[0].text = t
        btf.paragraphs[0].font.name = title_font
        btf.paragraphs[0].font.size = Pt(template.get('bodyFontSize', 16))
        if desc:
            p = btf.add_paragraph()
            p.text = desc
            p.font.name = body_font
            p.font.size = Pt(template.get('bodyFontSize', 14))
        y = y + Inches(1.6)
        if y > prs.slide_height - Inches(1):
            slide = prs.slides.add_slide(layout)
            y = Inches(0.5)
    notes = slide_def.get('notes')
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def create_education_slide(prs: Presentation, slide_def: Dict[str, Any], template: Dict[str, Any]):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    set_slide_background(prs, slide, template)
    title_font = _safe_font(slide, template.get('fontTitle'))
    body_font = _safe_font(slide, template.get('fontBody'))

    title = slide_def.get('title') or 'Education'
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), prs.slide_width - Inches(1), Inches(0.6)).text_frame
    tx.paragraphs[0].text = title
    tx.paragraphs[0].font.name = title_font
    tx.paragraphs[0].font.size = Pt(template.get('headingFontSize', 22))

    items = slide_def.get('items') or []
    y = Inches(1.0)
    for it in items:
        school = it.get('school') or ''
        degree = it.get('degree') or ''
        dates = it.get('dates') or ''
        box = slide.shapes.add_textbox(Inches(0.6), y, prs.slide_width - Inches(1.2), Inches(0.6))
        btf = box.text_frame
        btf.paragraphs[0].text = f"{degree} — {school} ({dates})"
        btf.paragraphs[0].font.name = title_font
        btf.paragraphs[0].font.size = Pt(template.get('bodyFontSize', 14))
        y = y + Inches(0.7)
    notes = slide_def.get('notes')
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def create_contact_slide(prs: Presentation, slide_def: Dict[str, Any], template: Dict[str, Any]):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    set_slide_background(prs, slide, template)
    title_font = _safe_font(slide, template.get('fontTitle'))
    body_font = _safe_font(slide, template.get('fontBody'))

    title = slide_def.get('title') or 'Contact'
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), prs.slide_width - Inches(1), Inches(0.6)).text_frame
    tx.paragraphs[0].text = title
    tx.paragraphs[0].font.name = title_font
    tx.paragraphs[0].font.size = Pt(template.get('headingFontSize', 22))

    bullets = slide_def.get('bullets') or []
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), prs.slide_width - Inches(1.2), Inches(4))
    btf = box.text_frame
    for b in bullets:
        p = btf.add_paragraph()
        p.text = b
        p.font.name = body_font
        p.font.size = Pt(template.get('bodyFontSize', 14))

    notes = slide_def.get('notes')
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


# Main builder

def build_presentation_from_deck(deck: Dict[str, Any]) -> io.BytesIO:
    # Use a base template if provided
    prs = Presentation()
    template = deck.get('template') or {}
    # If frontend didn't send a usable template, fall back to the built-in professional template
    try:
        if not template or not template.get('name'):
            import os, json
            base = os.path.join(os.path.dirname(__file__), 'templates')
            p = os.path.join(base, 'professional_clean.json')
            with open(p, 'r') as f:
                template = json.load(f)
    except Exception:
        # keep whatever template we have
        pass

    # Ensure some defaults exist for sizing/colors
    template.setdefault('accentColor', template.get('accentColor', '#0A74DA'))
    template.setdefault('bgColor', template.get('bgColor', '#FFFFFF'))
    template.setdefault('bgType', template.get('bgType', 'solid'))

    # iterate slides
    slides = deck.get('slides', []) or []
    # log basic info for debugging
    try:
        import logging
        logging.getLogger(__name__).info(f"Building PPTX: title={deck.get('title')} slides={len(slides)} template={template.get('name')}")
    except Exception:
        pass
    for s in slides:
        t = s.get('type')
        if t == 'title':
            create_title_slide(prs, s, template)
        elif t == 'summary':
            create_summary_slide(prs, s, template)
        elif t == 'experience':
            create_experience_slide(prs, s, template)
        elif t == 'skills':
            create_skills_slide(prs, s, template)
        elif t == 'projects':
            create_projects_slide(prs, s, template)
        elif t == 'education':
            create_education_slide(prs, s, template)
        elif t == 'contact':
            create_contact_slide(prs, s, template)
        else:
            # unknown slide -> create a simple text slide
            create_summary_slide(prs, s, template)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
