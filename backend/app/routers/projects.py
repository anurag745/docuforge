from fastapi import APIRouter, Depends, HTTPException
from typing import List
from app.schemas.schemas import ProjectCreate, ProjectOut, GenerateRequest, GenerateResponse, RefineRequest, FeedbackRequest, CommentRequest, ExportRequest, OutlineRequest, OutlineSuggestRequest
from app.database import get_session
from sqlalchemy.ext.asyncio import AsyncSession
from app.routers.auth import get_current_user
from app.controllers.project_controller import list_projects, create_project, get_project, append_revision, add_comment, log_llm, create_section, get_project_files
from app.models.models import Section
from app.services.llm_service import llm_service
from fastapi.responses import StreamingResponse
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from bs4 import BeautifulSoup
from typing import List
from fastapi import Body
from app.schemas.schemas import DeckModel, TemplateModel
from app.pptx_builder import build_presentation_from_deck
from fastapi import APIRouter

router = APIRouter()


async def get_db():
    async with get_session() as session:
        yield session


@router.get("/", response_model=List[ProjectOut])
async def projects(session: AsyncSession = Depends(get_db), user=Depends(get_current_user)):
    owner_id = int(user.get("sub"))
    res = await list_projects(session, owner_id)
    return res


@router.post("/", response_model=ProjectOut)
async def create(payload: ProjectCreate, session: AsyncSession = Depends(get_db), user=Depends(get_current_user)):
    owner_id = int(user.get("sub"))
    project = await create_project(session, owner_id, payload.title, payload.docType, payload.topic, payload.scaffold)
    return project


@router.get("/{project_id}", response_model=ProjectOut)
async def get_project_route(project_id: int, session: AsyncSession = Depends(get_db), user=Depends(get_current_user)):
    owner_id = int(user.get("sub"))
    project = await get_project(session, project_id, owner_id)
    return project


@router.delete("/{project_id}")
async def delete_project_route(project_id: int, session: AsyncSession = Depends(get_db), user=Depends(get_current_user)):
    owner_id = int(user.get("sub"))
    # call controller helper to perform delete
    from app.controllers.project_controller import delete_project as controller_delete
    await controller_delete(session, project_id, owner_id)
    return {"ok": True}


@router.post("/{project_id}/generate", response_model=GenerateResponse)
async def generate(project_id: int, payload: GenerateRequest, session: AsyncSession = Depends(get_db), user=Depends(get_current_user)):
    owner_id = int(user.get("sub"))
    # ensure project exists
    project = await get_project(session, project_id, owner_id)
    # determine section
    section_id = payload.sectionId
    section_obj = None
    if section_id:
        section_obj = await session.get(Section, section_id)
        if not section_obj:
            raise HTTPException(status_code=404, detail="Section not found")

    # provide context from section if available
    context = section_obj.content if section_obj and section_obj.content else None
    # pass project.docType and optional template so the LLM can return structured HTML for docx or slide HTML for pptx
    out = await llm_service.generate(project_id, section_id, payload.slideIndex, context=context, docType=project.docType, template=payload.template)

    # persist: if section exists, append revision; otherwise create a new section
    if section_obj:
        rev = await append_revision(session, section_obj.id, out["text"], payload.dict().__repr__())
    else:
        # create a new section with generated content
        title = f"Generated Section"
        sec = await create_section(session, project_id, title, out["text"], order_index=len(project.sections) if project.sections else 0)
        rev = await append_revision(session, sec.id, out["text"], payload.dict().__repr__())

    # log
    await log_llm(session, project_id, section_id, llm_service.provider, payload.dict().__repr__(), out["text"]) 
    return out


@router.post("/{project_id}/generate_all")
async def generate_all(project_id: int, session: AsyncSession = Depends(get_db), user=Depends(get_current_user)):
    owner_id = int(user.get("sub"))
    project = await get_project(session, project_id, owner_id)
    results = []
    # iterate section-by-section
    sections = project.sections or []
    for s in sections:
        out = await llm_service.generate(project_id, s.id, None, context=s.content, docType=project.docType)
        rev = await append_revision(session, s.id, out["text"], f"generate_all for section {s.id}")
        await log_llm(session, project_id, s.id, llm_service.provider, f"generate_all", out["text"])
        results.append({"sectionId": s.id, "text": out["text"], "generationId": out.get("generationId")})
    return {"results": results}


@router.post("/{project_id}/refine")
async def refine(project_id: int, payload: RefineRequest, session: AsyncSession = Depends(get_db), user=Depends(get_current_user)):
    owner_id = int(user.get("sub"))
    _ = await get_project(session, project_id, owner_id)
    # get section content
    s = await session.get(Section, payload.sectionId)
    current_text = s.content if s and s.content else ""
    refined = await llm_service.refine(current_text, payload.prompt)
    rev = await append_revision(session, payload.sectionId, refined, payload.prompt)
    await log_llm(session, project_id, payload.sectionId, llm_service.provider, payload.prompt, refined)
    return {"text": refined, "revisionId": rev.id}


@router.post("/{project_id}/feedback")
async def feedback(project_id: int, payload: FeedbackRequest, session: AsyncSession = Depends(get_db), user=Depends(get_current_user)):
    owner_id = int(user.get("sub"))
    _ = await get_project(session, project_id, owner_id)
    # store as a log for now
    await log_llm(session, project_id, payload.sectionId, llm_service.provider, f"feedback: {payload.like}", "")
    return {"ok": True}


@router.post("/{project_id}/comment")
async def comment(project_id: int, payload: CommentRequest, session: AsyncSession = Depends(get_db), user=Depends(get_current_user)):
    owner_id = int(user.get("sub"))
    _ = await get_project(session, project_id, owner_id)
    c = await add_comment(session, payload.sectionId, int(owner_id), payload.comment)
    return c


@router.post("/{project_id}/suggest_outline")
async def suggest_outline(project_id: int, payload: OutlineRequest, session: AsyncSession = Depends(get_db), user=Depends(get_current_user)):
    owner_id = int(user.get("sub"))
    project = await get_project(session, project_id, owner_id)
    titles = await llm_service.suggest_outline(payload.topic, project.docType, template=payload.template)
    return {"titles": titles}


@router.post("/suggest_outline")
async def suggest_outline_no_project(payload: OutlineSuggestRequest, session: AsyncSession = Depends(get_db), user=Depends(get_current_user)):
    owner_id = int(user.get("sub"))
    titles = await llm_service.suggest_outline(payload.topic, payload.docType, template=payload.template)
    return {"titles": titles}


@router.get("/{project_id}/files")
async def get_files(project_id: int, session: AsyncSession = Depends(get_db), user=Depends(get_current_user)):
    owner_id = int(user.get("sub"))
    files = await get_project_files(session, project_id, owner_id)
    # Pydantic will serialize ORM objects when returned by FastAPI if configured; otherwise return raw dict-like
    return files


@router.post("/{project_id}/export")
async def export(project_id: int, payload: ExportRequest, session: AsyncSession = Depends(get_db), user=Depends(get_current_user)):
    owner_id = int(user.get("sub"))
    project = await get_project(session, project_id, owner_id)
    fmt = (payload.format or "txt").lower()

    # If pptx requested, build a real PowerPoint file from sections
    if fmt == "pptx":
        prs = Presentation()

        # Choose a blank slide layout (0 is usually title slide, 5/6 blank/content - safer to use layout 5 if available)
        # We'll attempt to use layout 1 (Title and Content) if present
        layout_index = 1 if len(prs.slide_layouts) > 1 else 0

        def html_to_slide_parts(title: str, html: str) -> (str, List[str]):
            """Parse html fragment into a title and bullet/paragraph lines."""
            if not html:
                return title or "", []
            try:
                soup = BeautifulSoup(html, "html.parser")
                # prefer heading tags inside html
                h = None
                for tag in ("h1", "h2", "h3", "h4"):
                    el = soup.find(tag)
                    if el and el.get_text(strip=True):
                        h = el.get_text(strip=True)
                        break
                slide_title = h or title or ""

                # collect bullet items from ul/li, otherwise paragraphs
                bullets = []
                for ul in soup.find_all("ul"):
                    for li in ul.find_all("li"):
                        text = li.get_text(" ", strip=True)
                        if text:
                            bullets.append(text)
                if not bullets:
                    # fallback to paragraphs
                    for p in soup.find_all("p"):
                        text = p.get_text(" ", strip=True)
                        if text:
                            bullets.append(text)

                # if still empty, use the whole text
                if not bullets:
                    all_text = soup.get_text(" ", strip=True)
                    if all_text:
                        bullets = [all_text]

                return slide_title, bullets
            except Exception:
                # fallback: strip tags roughly
                stripped = re.sub(r"<[^>]+>", "", html)
                return title or "", [stripped]

        # Iterate sections and add slides
        for s in project.sections or []:
            title_text, lines = html_to_slide_parts(s.title, s.content or "")
            slide_layout = prs.slide_layouts[layout_index]
            slide = prs.slides.add_slide(slide_layout)

            # Try to set title if the layout has a title placeholder
            try:
                title_shape = slide.shapes.title
                if title_shape and title_text:
                    title_shape.text = title_text
            except Exception:
                # no title placeholder, skip
                pass

            # Find a content placeholder and add bullets/paragraphs
            body = None
            for shape in slide.shapes:
                if not shape.is_placeholder:
                    continue
                phf = shape.placeholder_format
                # content placeholders often have idx 1
                try:
                    if phf.type.name in ("BODY", "CONTENT"):
                        body = shape
                        break
                except Exception:
                    continue

            if body is None:
                # create a textbox
                left = Inches(1)
                top = Inches(1.8)
                width = Inches(8)
                height = Inches(4.5)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
            else:
                tf = body.text_frame

            # clear default paragraph
            try:
                tf.clear()
            except Exception:
                pass

            for i, line in enumerate(lines):
                if i == 0:
                    p = tf.paragraphs[0]
                    p.text = line
                else:
                    p = tf.add_paragraph()
                    p.text = line
                    p.level = 1
                # styling
                p.font.size = Pt(18)

        # save presentation to buffer
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        filename = f"project-{project_id}.pptx"
        return StreamingResponse(buf, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", headers={"Content-Disposition": f"attachment; filename=\"{filename}\""})

    # fallback: build a simple text file as mock
    buf = io.BytesIO()
    content = f"Export for project {project.title}\nFormat: {payload.format}\nSections: {payload.sections}\nIncludeComments: {payload.includeComments}\n"
    buf.write(content.encode("utf-8"))
    buf.seek(0)
    filename = f"project-{project_id}." + (payload.format or "txt")
    return StreamingResponse(buf, media_type="application/octet-stream", headers={"Content-Disposition": f"attachment; filename=\"{filename}\""})


@router.post('/suggest_templates')
async def suggest_templates(payload: dict = Body(...)):
    """Return built-in templates and (optionally) OpenAI-suggested templates.

    Body: { "topic": "job role or candidate summary" }
    Currently returns the three built-in templates from `app/templates/`.
    """
    import json, os
    base = os.path.join(os.path.dirname(__file__), '..', 'templates')
    templates = []
    try:
        for fname in ('professional_clean.json','modern_minimal.json','creative_portfolio.json'):
            p = os.path.join(base, fname)
            with open(p, 'r') as f:
                templates.append(json.load(f))
    except Exception:
        templates = []
    return {"templates": templates}


@router.post('/generate_pptx')
async def generate_pptx(deck: DeckModel, user=Depends(get_current_user)):
    """Generate a PPTX from DeckModel and stream it back.

    The DeckModel is validated by Pydantic.
    """
    # Build presentation
    buf = build_presentation_from_deck(deck.dict())
    safe_title = re.sub(r"[^0-9A-Za-z-_]+", "_", deck.title) if deck.title else 'presentation'
    filename = f"{safe_title}.pptx"
    return StreamingResponse(buf, media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation', headers={"Content-Disposition": f"attachment; filename=\"{filename}\""})
